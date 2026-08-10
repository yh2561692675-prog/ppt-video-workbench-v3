from __future__ import annotations

from pathlib import Path

from docx import Document


def test_p04_extracts_docx_outline_and_page_payload(tmp_path: Path) -> None:
    from workbench.business_modules.p04_extract.runner import extract_document

    source = tmp_path / "outline.docx"
    document = Document()
    document.add_heading("课程概览", level=1)
    document.add_paragraph("这是课程介绍。")
    document.save(source)

    payload = extract_document(source, tmp_path / "previews", "auto")

    assert payload["outline"]["blocks"][0]["text"] == "课程概览"
    assert payload["page_count"] == 0
    assert payload["pages"] == []


def test_p04_extracts_pptx_slides(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    from workbench.business_modules.p04_extract.runner import extract_document

    source = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "第一页"
    presentation.save(source)

    payload = extract_document(source, tmp_path / "previews", "never")

    assert payload["page_count"] == 1
    assert payload["pages"][0]["title"] == "第一页"
