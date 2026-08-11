from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches
from workbench.fidelity.models import FidelityLevel, FidelityPolicy
from workbench.fidelity.scanner import FidelityScanError, PptxFidelityScanner


def _pptx(path: Path, slide_count: int = 1) -> None:
    presentation = Presentation()
    for index in range(slide_count):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        shape.text = f"Fidelity page {index + 1}"
    presentation.save(path)


def test_scanner_extracts_stable_shape_scene_and_hash(tmp_path: Path) -> None:
    source = tmp_path / "deck.pptx"
    _pptx(source, slide_count=2)

    first = PptxFidelityScanner().scan(source, tmp_path / "fidelity")
    second = PptxFidelityScanner().scan(source, tmp_path / "fidelity-2")

    assert first.pages[0].level is FidelityLevel.F0
    assert first.pages[0].scene.shapes[0].text == "Fidelity page 1"
    assert first.source_hash == second.source_hash
    assert first.pages[0].page_id == second.pages[0].page_id
    assert first.manifest_hash == second.manifest_hash


def test_scanner_uses_static_renderer_and_enforces_limits(tmp_path: Path) -> None:
    source = tmp_path / "deck.pptx"
    _pptx(source, slide_count=2)
    preview = tmp_path / "fidelity" / "page-1.png"

    def renderer(_source: Path, output: Path) -> list[Path]:
        preview.parent.mkdir(parents=True, exist_ok=True)
        preview.write_bytes(b"png")
        return [preview]

    manifest = PptxFidelityScanner(static_renderer=renderer).scan(source, tmp_path / "fidelity")
    assert manifest.pages[0].level is FidelityLevel.F1
    assert manifest.pages[0].preview_path == preview.as_posix()
    with pytest.raises(FidelityScanError):
        PptxFidelityScanner().scan(source, tmp_path / "limited", FidelityPolicy(max_slide_count=1))


def test_scanner_rejects_active_content_and_zip_path_escape(tmp_path: Path) -> None:
    active = tmp_path / "active.pptx"
    with zipfile.ZipFile(active, "w") as archive:
        archive.writestr("ppt/embeddings/oleObject1.bin", b"bad")
    with pytest.raises(FidelityScanError, match="宏"):
        PptxFidelityScanner().scan(active, tmp_path / "out")

    escaped = tmp_path / "escaped.pptx"
    with zipfile.ZipFile(escaped, "w") as archive:
        archive.writestr("../escape.txt", b"bad")
    with pytest.raises(FidelityScanError, match="越界"):
        PptxFidelityScanner().scan(escaped, tmp_path / "out")
