from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from workbench.parsers.pptx_parser import parse_pptx
from workbench.renderers.office_renderer import (
    MissingFontError,
    OfficeRendererError,
    OfficeRendererUnavailable,
    build_pptx_previews,
    ensure_fonts_available,
    render_office_to_pdf,
)


def build_presentation(path: Path, *, extra_pages: int = 0) -> None:
    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    first = deck.slides.add_slide(deck.slide_layouts[6])
    box = first.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text = "第一页标题"
    box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
    table = first.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "方向"
    table.cell(0, 1).text = "课程"
    table.cell(1, 0).text = "人工智能"
    table.cell(1, 1).text = "机器学习"
    first.notes_slide.notes_text_frame.text = "备注秘密，不应提取"

    second = deck.slides.add_slide(deck.slide_layouts[6])
    second_box = second.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    second_box.text = "第二页标题"
    second_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
    accent = second.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(20, 90, 140)
    second._element.set("show", "0")
    for index in range(extra_pages):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        title.text = f"专业介绍第 {index + 3} 页"
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(20 + index * 10, 90, 140)
    deck.save(path)


def test_parse_pptx_preserves_page_order_tables_and_hidden_flag(tmp_path: Path) -> None:
    path = tmp_path / "课件.pptx"
    build_presentation(path)

    pages = parse_pptx(path)

    assert [page.order for page in pages] == [1, 2]
    assert pages[0].text == "第一页标题\n方向\n课程\n人工智能\n机器学习"
    assert "备注秘密" not in pages[0].text
    assert pages[0].hidden is False
    assert pages[1].hidden is True


def test_office_render_fails_closed_without_engine_or_required_font(tmp_path: Path) -> None:
    path = tmp_path / "课件.pptx"
    build_presentation(path)

    with pytest.raises(OfficeRendererUnavailable, match="LibreOffice"):
        render_office_to_pdf(path, tmp_path / "pdf", executable=tmp_path / "missing-soffice")
    with pytest.raises(MissingFontError, match="缺少字体"):
        ensure_fonts_available({"Definitely Missing Chinese Font"})

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text = "缺失字体"
    box.text_frame.paragraphs[0].runs[0].font.name = "Definitely Missing Chinese Font"
    deck.save(path)
    with pytest.raises(MissingFontError, match="缺少字体"):
        build_pptx_previews(path, tmp_path / "missing-font-preview")


def test_office_render_never_accepts_stale_pdf_after_failed_conversion(tmp_path: Path) -> None:
    path = tmp_path / "课件.pptx"
    build_presentation(path)
    output = tmp_path / "pdf"
    output.mkdir()
    (output / "课件.pdf").write_bytes(b"%PDF-stale")
    fake = tmp_path / "fake-soffice"
    fake.write_text("#!/bin/sh\nexit 0\n", encoding="utf-8")
    fake.chmod(0o755)

    with pytest.raises(OfficeRendererError, match="未生成 PDF"):
        render_office_to_pdf(path, output, executable=fake)


def test_real_pptx_render_produces_nonblank_full_hd_pages(tmp_path: Path) -> None:
    path = tmp_path / "课件.pptx"
    build_presentation(path, extra_pages=6)

    result = build_pptx_previews(path, tmp_path / "预览")

    assert result.engine.startswith("LibreOffice ")
    assert len(result.pages) == 8
    assert result.pages[7].text == "专业介绍第 8 页"
    for page in result.pages:
        with Image.open(page.preview_path) as preview:
            assert preview.size == (1920, 1080)
            assert preview.getbbox() is not None
