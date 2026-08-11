from __future__ import annotations

from pathlib import Path
from uuid import uuid4

from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from workbench.api.fidelity import create_fidelity_router
from workbench.fidelity.jobs import FidelityJobService
from workbench.fidelity.scanner import PptxFidelityScanner


def test_fidelity_routes_create_job_and_list_pages(tmp_path: Path) -> None:
    project_id = uuid4()
    project_root = tmp_path / str(project_id)
    project_root.mkdir()
    source = project_root / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    shape.text = "页面"
    presentation.save(source)
    app = FastAPI()
    app.include_router(create_fidelity_router(FidelityJobService(tmp_path, PptxFidelityScanner())))

    with TestClient(app) as client:
        response = client.post(
            f"/api/projects/{project_id}/fidelity/jobs",
            json={"pptx_path": "deck.pptx"},
        )
        assert response.status_code == 201
        assert response.json()["data"]["status"] == "succeeded"
        pages = client.get(f"/api/projects/{project_id}/fidelity/pages")
        assert pages.status_code == 200
        assert pages.json()["data"][0]["page_index"] == 1
