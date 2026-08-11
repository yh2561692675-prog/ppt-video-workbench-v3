from __future__ import annotations

import argparse
import math
import struct
import wave
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

SLIDES = (
    ("Windows 全流程验收", "安装、启动与项目导入", RGBColor(30, 64, 175)),
    ("单页渲染", "验证页面级预览、声音和画面", RGBColor(5, 150, 105)),
    ("批量渲染与恢复", "模拟中断后继续执行未完成页面", RGBColor(217, 119, 6)),
    ("最终合成", "导出视频并验证重启、回滚和卸载", RGBColor(190, 24, 93)),
)


def create_deck(path: Path) -> None:
    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    for index, (title, subtitle, accent) in enumerate(SLIDES, start=1):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.45), Inches(10.8), Inches(1.2))
        title_run = title_box.text_frame.paragraphs[0].add_run()
        title_run.text = title
        title_run.font.name = "Microsoft YaHei"
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(15, 23, 42)

        subtitle_box = slide.shapes.add_textbox(Inches(1.1), Inches(3.0), Inches(10.8), Inches(0.8))
        subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
        subtitle_run.text = subtitle
        subtitle_run.font.name = "Microsoft YaHei"
        subtitle_run.font.size = Pt(20)
        subtitle_run.font.color.rgb = RGBColor(71, 85, 105)

        marker = slide.shapes.add_textbox(Inches(11.35), Inches(6.45), Inches(1.0), Inches(0.5))
        marker_run = marker.text_frame.paragraphs[0].add_run()
        marker_run.text = f"{index:02d}"
        marker_run.font.name = "Segoe UI"
        marker_run.font.size = Pt(16)
        marker_run.font.bold = True
        marker_run.font.color.rgb = accent

    path.parent.mkdir(parents=True, exist_ok=True)
    deck.save(path)


def create_outline(path: Path) -> None:
    document = Document()
    document.add_heading("Windows 全流程验收", level=0)
    for title, subtitle, _ in SLIDES:
        document.add_heading(title, level=1)
        document.add_paragraph(subtitle)
        document.add_paragraph(f"{title}的验收旁白。{subtitle}。")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def create_audio(path: Path, *, duration_seconds: float = 16.0) -> None:
    sample_rate = 48_000
    amplitude = 7_000
    frame_count = int(sample_rate * duration_seconds)
    path.parent.mkdir(parents=True, exist_ok=True)
    with wave.open(str(path), "wb") as handle:
        handle.setnchannels(1)
        handle.setsampwidth(2)
        handle.setframerate(sample_rate)
        frames = bytearray()
        for index in range(frame_count):
            second = index / sample_rate
            section = min(int(second // 4), len(SLIDES) - 1)
            carrier = 180 + section * 45
            envelope = 0.45 + 0.45 * math.sin(2 * math.pi * 1.7 * second) ** 2
            sample = int(
                amplitude
                * envelope
                * (
                    math.sin(2 * math.pi * carrier * second)
                    + 0.35 * math.sin(2 * math.pi * carrier * 1.5 * second)
                )
            )
            frames.extend(struct.pack("<h", max(-32768, min(32767, sample))))
        handle.writeframes(frames)


def main() -> int:
    parser = argparse.ArgumentParser(description="Create deterministic Windows full-flow fixtures")
    parser.add_argument("--output-dir", type=Path, required=True)
    args = parser.parse_args()
    output_dir = args.output_dir.resolve()
    create_deck(output_dir / "windows-full-flow.pptx")
    create_outline(output_dir / "windows-full-flow-outline.docx")
    create_audio(output_dir / "windows-full-flow-narration.wav")
    print(f"FIXTURE_DECK={output_dir / 'windows-full-flow.pptx'}")
    print(f"FIXTURE_OUTLINE={output_dir / 'windows-full-flow-outline.docx'}")
    print(f"FIXTURE_AUDIO={output_dir / 'windows-full-flow-narration.wav'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
