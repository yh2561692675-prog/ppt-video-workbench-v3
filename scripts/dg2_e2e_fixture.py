"""Create and validate deterministic, synthetic DG2 browser fixtures.

The generated files deliberately live under ``tests/.e2e-fixtures`` (which is
ignored).  The recipe and the expected SHA-256 values are versioned instead of
checking a private recording or a machine-generated binary into the repository.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import shutil
import struct
import tempfile
import wave
import zipfile
from datetime import UTC, datetime
from pathlib import Path

from docx import Document
from pptx import Presentation

FIXTURE_VERSION = "dg2-synthetic-v1"
FIXED_TIME = datetime(2026, 8, 11, tzinfo=UTC)
ROOT = Path(__file__).resolve().parents[1]
CONTRACT_PATH = ROOT / "fixtures" / "dg2" / "fixture-contract-v1.json"

NARRATIONS = tuple(f"合成第{number}页旁白" for number in range(1, 9))
PROFILES = {"S1": 2, "S8": 8}


class FixtureValidationError(ValueError):
    """Raised when a generated fixture no longer matches its reviewed contract."""


def fixture_directory(output_root: Path, profile: str) -> Path:
    if profile not in PROFILES:
        raise FixtureValidationError(f"unsupported DG2 fixture profile: {profile}")
    return output_root / profile.lower()


def generate_fixtures(output_root: Path) -> dict[str, dict[str, object]]:
    """Generate the S1 and S8 sources and return their observed manifests."""

    output_root.mkdir(parents=True, exist_ok=True)
    result: dict[str, dict[str, object]] = {}
    for profile, page_count in PROFILES.items():
        target = fixture_directory(output_root, profile)
        if target.exists():
            shutil.rmtree(target)
        target.mkdir(parents=True)
        _create_docx(target / "outline.docx", page_count)
        _create_pptx(target / "deck.pptx", page_count)
        _create_wav(target / "local-narration.wav", page_count)
        result[profile] = _manifest(profile, target, page_count)
        (target / "fixture-manifest.json").write_text(
            json.dumps(result[profile], ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
    return result


def validate_fixtures(
    output_root: Path,
    *,
    contract_path: Path = CONTRACT_PATH,
) -> dict[str, dict[str, object]]:
    """Validate generated sources against fixed content, shape and SHA-256 values."""

    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    if contract.get("fixture_version") != FIXTURE_VERSION:
        raise FixtureValidationError("DG2 fixture contract version does not match the generator")
    observed: dict[str, dict[str, object]] = {}
    for profile, page_count in PROFILES.items():
        folder = fixture_directory(output_root, profile)
        manifest = _manifest(profile, folder, page_count)
        expected = contract.get("profiles", {}).get(profile)
        if not isinstance(expected, dict):
            raise FixtureValidationError(f"missing contract for {profile}")
        _validate_profile(profile, manifest, expected)
        observed[profile] = manifest
    return observed


def _create_docx(path: Path, page_count: int) -> None:
    document = Document()
    properties = document.core_properties
    properties.author = "PPT Video Workbench DG2"
    properties.title = f"{page_count} 页合成大纲"
    properties.subject = FIXTURE_VERSION
    properties.created = FIXED_TIME
    properties.modified = FIXED_TIME
    for number in range(1, page_count + 1):
        document.add_heading(f"合成第{number}页", level=1)
        document.add_paragraph(NARRATIONS[number - 1])
    document.save(str(path))
    _normalise_zip(path)


def _create_pptx(path: Path, page_count: int) -> None:
    presentation = Presentation()
    properties = presentation.core_properties
    properties.author = "PPT Video Workbench DG2"
    properties.title = f"{page_count} 页合成课件"
    properties.subject = FIXTURE_VERSION
    properties.created = FIXED_TIME
    properties.modified = FIXED_TIME
    for number in range(1, page_count + 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"合成第{number}页"
        slide.placeholders[1].text = (
            f"DG2 固定合成课件，第 {number} 页。\n{NARRATIONS[number - 1]}。"
        )
    presentation.save(str(path))
    _normalise_zip(path)


def _create_wav(path: Path, page_count: int) -> None:
    sample_rate = 16_000
    duration_ms = page_count * 750
    frames = duration_ms * sample_rate // 1_000
    with wave.open(str(path), "wb") as handle:
        handle.setnchannels(1)
        handle.setsampwidth(2)
        handle.setframerate(sample_rate)
        samples = bytearray()
        for index in range(frames):
            value = int(4_000 * math.sin(2 * math.pi * 440 * index / sample_rate))
            samples.extend(struct.pack("<h", value))
        handle.writeframes(bytes(samples))


def _normalise_zip(path: Path) -> None:
    """Canonicalise OOXML ZIP metadata so the reviewed hash is cross-run stable."""

    with tempfile.TemporaryDirectory(prefix="dg2-ooxml-") as temporary:
        canonical = Path(temporary) / path.name
        with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(
            canonical,
            "w",
            compression=zipfile.ZIP_STORED,
        ) as target:
            for name in sorted(source.namelist()):
                info = zipfile.ZipInfo(filename=name, date_time=(1980, 1, 1, 0, 0, 0))
                # Stored entries avoid zlib implementation differences between
                # Windows and Linux changing the reviewed fixture hashes.
                info.compress_type = zipfile.ZIP_STORED
                info.external_attr = 0o600 << 16
                target.writestr(info, source.read(name), compress_type=zipfile.ZIP_STORED)
        shutil.copy2(canonical, path)


def _manifest(profile: str, folder: Path, page_count: int) -> dict[str, object]:
    files = {}
    for name in ("outline.docx", "deck.pptx", "local-narration.wav"):
        path = folder / name
        if not path.is_file():
            raise FixtureValidationError(f"{profile} is missing {name}")
        files[name] = {"sha256": _content_sha256(path), "size": path.stat().st_size}
    with wave.open(str(folder / "local-narration.wav"), "rb") as handle:
        audio = {
            "sample_rate": handle.getframerate(),
            "channels": handle.getnchannels(),
            "duration_ms": round(handle.getnframes() * 1_000 / handle.getframerate()),
        }
    deck = Presentation(str(folder / "deck.pptx"))
    if len(deck.slides) != page_count:
        raise FixtureValidationError(f"{profile} PPTX has unexpected slide count")
    return {
        "fixture_version": FIXTURE_VERSION,
        "profile": profile,
        "content_policy": "synthetic-only",
        "page_count": page_count,
        "narrations": list(NARRATIONS[:page_count]),
        "audio": audio,
        "files": files,
        "expected_outputs": [
            "08_输出/最终视频.mp4",
            "08_输出/制作包/字幕.srt",
            "08_输出/制作包/制作包清单.json",
        ],
    }


def _validate_profile(
    profile: str,
    observed: dict[str, object],
    expected: dict[str, object],
) -> None:
    for key in ("content_policy", "page_count", "narrations", "audio", "expected_outputs"):
        if observed.get(key) != expected.get(key):
            raise FixtureValidationError(
                f"{profile} fixture {key} drifted from the reviewed contract"
            )
    observed_files = observed.get("files")
    expected_files = expected.get("files")
    if observed_files != expected_files:
        raise FixtureValidationError(
            f"{profile} fixture source hashes drifted from the reviewed contract"
        )


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        for chunk in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _content_sha256(path: Path) -> str:
    """Hash OOXML package contents independently of platform-specific ZIP bytes."""

    if path.suffix.lower() not in {".docx", ".pptx"}:
        return _sha256(path)
    digest = hashlib.sha256()
    with zipfile.ZipFile(path, "r") as package:
        for name in sorted(package.namelist()):
            name_bytes = name.encode("utf-8")
            content = package.read(name)
            digest.update(struct.pack("<I", len(name_bytes)))
            digest.update(name_bytes)
            digest.update(struct.pack("<Q", len(content)))
            digest.update(content)
    return digest.hexdigest()


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument(
        "--bootstrap",
        action="store_true",
        help="print observed fixture contracts without requiring reviewed hash values",
    )
    args = parser.parse_args()
    observed = generate_fixtures(args.output)
    if args.bootstrap:
        print(
            json.dumps(
                {"fixture_version": FIXTURE_VERSION, "profiles": observed},
                ensure_ascii=False,
            )
        )
        return 0
    validate_fixtures(args.output)
    print(json.dumps({"status": "passed", "fixture_version": FIXTURE_VERSION}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
