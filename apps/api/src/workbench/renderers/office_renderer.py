from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

import fitz  # type: ignore[import-untyped]
from PIL import Image
from pptx import Presentation

from workbench.domain.extraction import PageExtraction, PreviewBuildResult
from workbench.parsers.pptx_parser import parse_pptx


class OfficeRendererError(RuntimeError):
    pass


class OfficeRendererUnavailable(OfficeRendererError):
    pass


class MissingFontError(OfficeRendererError):
    pass


def ensure_fonts_available(required_fonts: set[str]) -> None:
    matcher = shutil.which("fc-match")
    if not matcher:
        if required_fonts:
            raise MissingFontError(
                "\u7f3a\u5c11\u5b57\u4f53\uff1a" + ", ".join(sorted(required_fonts))
            )
        return
    missing: list[str] = []
    for font in sorted(required_fonts):
        result = subprocess.run(
            [matcher, "-f", "%{family}", font], capture_output=True, text=True, check=False
        )
        families = result.stdout.casefold()
        if font.casefold() not in families:
            missing.append(font)
    if missing:
        raise MissingFontError(f"缺少字体：{', '.join(missing)}")


def render_office_to_pdf(path: Path, out: Path, *, executable: Path | None = None) -> Path:
    binary = executable or _office_binary()
    if not binary.exists():
        raise OfficeRendererUnavailable("未找到 LibreOffice，无法生成页面预览")
    out.mkdir(parents=True, exist_ok=True)
    target = out / f"{path.stem}.pdf"
    target.unlink(missing_ok=True)
    profile = out / ".libreoffice-profile"
    profile.mkdir(exist_ok=True)
    command = [
        str(binary),
        "--headless",
        f"-env:UserInstallation={profile.resolve().as_uri()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out),
        str(path),
    ]
    try:
        result = subprocess.run(command, capture_output=True, text=True, timeout=120, check=False)
    except (OSError, subprocess.TimeoutExpired) as error:
        raise OfficeRendererError("LibreOffice 转换失败：未生成 PDF") from error
    if result.returncode != 0 or not target.exists() or target.stat().st_size == 0:
        detail = (result.stderr or result.stdout).strip()
        raise OfficeRendererError(f"LibreOffice 转换失败：{detail or '未生成 PDF'}")
    return target


def build_pptx_previews(path: Path, out: Path) -> PreviewBuildResult:
    ensure_fonts_available(_pptx_fonts(path))
    pages = parse_pptx(path)
    render_source = _visible_render_copy(path, out / "_office")
    pdf = render_office_to_pdf(render_source, out / "_office")
    document = fitz.open(pdf)
    if len(document) != len(pages):
        document.close()
        raise OfficeRendererError("转换页数与原 PPTX 不一致")
    out.mkdir(parents=True, exist_ok=True)
    rendered: list[PageExtraction] = []
    try:
        for extraction, pdf_page in zip(pages, document, strict=True):
            scale = min(1920 / pdf_page.rect.width, 1080 / pdf_page.rect.height)
            pixmap = pdf_page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
            image = pixmap.pil_image().convert("RGB")
            canvas = Image.new("RGB", (1920, 1080), "white")
            offset = ((1920 - image.width) // 2, (1080 - image.height) // 2)
            canvas.paste(image, offset)
            colors = canvas.getcolors(maxcolors=2)
            if colors is not None and len(colors) == 1:
                raise OfficeRendererError(f"第 {extraction.order} 页渲染为空白页")
            preview = out / f"page-{extraction.order:04d}.png"
            canvas.save(preview, format="PNG")
            rendered.append(
                extraction.model_copy(
                    update={"preview_path": preview, "width": 1920, "height": 1080}
                )
            )
    finally:
        document.close()
    return PreviewBuildResult(engine=_office_version(), pages=rendered)


def _visible_render_copy(path: Path, cache_dir: Path) -> Path:
    presentation = Presentation(str(path))
    if not any(slide._element.get("show") == "0" for slide in presentation.slides):
        return path
    cache_dir.mkdir(parents=True, exist_ok=True)
    target = cache_dir / f"{path.stem}-all-pages.pptx"
    for slide in presentation.slides:
        if slide._element.get("show") == "0":
            del slide._element.attrib["show"]
    presentation.save(str(target))
    return target


def _pptx_fonts(path: Path) -> set[str]:
    presentation = Presentation(str(path))
    fonts: set[str] = set()
    for slide in presentation.slides:
        for shape in slide.shapes:
            text_frames = []
            if getattr(shape, "has_text_frame", False):
                text_frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                text_frames.extend(
                    cell.text_frame for row in shape.table.rows for cell in row.cells
                )
            for text_frame in text_frames:
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and not run.font.name.startswith("+"):
                            fonts.add(run.font.name)
    return fonts


def _office_binary() -> Path:
    located = shutil.which("soffice") or shutil.which("libreoffice")
    if not located:
        raise OfficeRendererUnavailable("未找到 LibreOffice，无法生成页面预览")
    return Path(located)


def _office_version() -> str:
    binary = _office_binary()
    result = subprocess.run([str(binary), "--version"], capture_output=True, text=True, check=False)
    version = result.stdout.strip() or result.stderr.strip() or "unknown"
    return version if version.startswith("LibreOffice ") else f"LibreOffice {version}"
