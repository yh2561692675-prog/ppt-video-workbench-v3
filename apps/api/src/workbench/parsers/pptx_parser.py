from __future__ import annotations

import hashlib
from pathlib import Path
from uuid import NAMESPACE_URL, uuid5

from pptx import Presentation

from workbench.domain.extraction import PageExtraction


class PresentationParseError(ValueError):
    pass


def parse_pptx(path: Path) -> list[PageExtraction]:
    try:
        presentation = Presentation(str(path))
    except Exception as error:
        raise PresentationParseError(f"无法解析 PPTX：{path.name}") from error
    source_hash = hashlib.sha256(path.read_bytes()).hexdigest()
    pages: list[PageExtraction] = []
    for order, slide in enumerate(presentation.slides, start=1):
        fragments: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    fragments.extend(cell.text.strip() for cell in row.cells if cell.text.strip())
            elif getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    fragments.extend(line.strip() for line in text.splitlines() if line.strip())
        pages.append(
            PageExtraction(
                id=uuid5(NAMESPACE_URL, f"pptx:{source_hash}:{order}"),
                order=order,
                text="\n".join(fragments),
                title=fragments[0] if fragments else None,
                hidden=slide._element.get("show") == "0",
                extraction_method="pptx",
                source_ref=f"slide:{order}",
            )
        )
    return pages
