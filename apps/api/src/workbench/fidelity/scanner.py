from __future__ import annotations

import hashlib
import json
import re
import zipfile
from collections.abc import Callable
from pathlib import Path
from typing import Any
from uuid import NAMESPACE_URL, UUID, uuid5

from pptx import Presentation

from .models import (
    FidelityIssue,
    FidelityLevel,
    FidelityPolicy,
    FidelityRenderer,
    MotionCue,
    MotionSupport,
    SlideFidelityManifest,
    SlideFidelityPage,
    SlideScene,
    SlideShape,
)


class FidelityScanError(RuntimeError):
    def __init__(self, code: str, message: str, action: str = "检查 PPTX 后重试") -> None:
        super().__init__(message)
        self.code = code
        self.action = action


StaticRenderer = Callable[[Path, Path], list[Path]]


class PptxFidelityScanner:
    def __init__(
        self,
        *,
        static_renderer: StaticRenderer | None = None,
        scanner_version: str = "fidelity-scanner-v1",
    ) -> None:
        self.static_renderer = static_renderer
        self.scanner_version = scanner_version

    def scan(
        self, path: Path, output_dir: Path, policy: FidelityPolicy | None = None
    ) -> SlideFidelityManifest:
        selected = policy or FidelityPolicy()
        path = path.resolve()
        if not path.is_file() or path.suffix.lower() != ".pptx":
            raise FidelityScanError("fidelity_input_invalid", "输入文件必须是存在的 PPTX 文件")
        source_hash = _sha256(path)
        self._validate_package(path, selected)
        try:
            presentation = Presentation(str(path))
        except Exception as error:  # noqa: BLE001 - parser boundary
            raise FidelityScanError("fidelity_parse_failed", "PPTX 无法安全解析") from error
        if len(presentation.slides) > selected.max_slide_count:
            raise FidelityScanError("fidelity_slide_count_limit", "PPTX 页数超过安全上限")
        output_dir.mkdir(parents=True, exist_ok=True)
        previews = self.static_renderer(path, output_dir) if self.static_renderer else []
        pages: list[SlideFidelityPage] = []
        for index, slide in enumerate(presentation.slides, start=1):
            page_id = uuid5(NAMESPACE_URL, f"pptx:{source_hash}:{index}")
            scene, issues = _extract_scene(
                slide,
                index,
                page_id,
                int(presentation.slide_width or 1),
                int(presentation.slide_height or 1),
            )
            motion = _extract_motion(path, index, page_id)
            scene = scene.model_copy(update={"motion_cues": motion})
            if motion and any(cue.support is not MotionSupport.SUPPORTED for cue in motion):
                issues.append(
                    FidelityIssue(
                        code="animation_degraded",
                        message="部分 PowerPoint 动画无法稳定映射为自动视频动画",
                        action="使用静态降级或启用 PowerPoint 原生捕获",
                    )
                )
            renderer = FidelityRenderer.LIBREOFFICE if previews else FidelityRenderer.PYTHON
            level = FidelityLevel.F1 if previews else FidelityLevel.F0
            if motion:
                level = FidelityLevel.F2
            if selected.require_animation_support and any(
                cue.support is not MotionSupport.SUPPORTED for cue in motion
            ):
                level = FidelityLevel.F0
            preview_path = previews[index - 1].as_posix() if index <= len(previews) else None
            pages.append(
                SlideFidelityPage(
                    page_id=page_id,
                    page_index=index,
                    level=level,
                    renderer=renderer,
                    scene=scene,
                    preview_path=preview_path,
                    issues=issues,
                    downgrade_reason=(
                        "animation_mapping_unsupported"
                        if level is FidelityLevel.F0 and motion
                        else None
                    ),
                    input_hash=source_hash,
                )
            )
        manifest = SlideFidelityManifest(
            source_path=path.as_posix(),
            source_hash=source_hash,
            pages=pages,
            capability={
                "pptx_parse": True,
                "static_render": bool(previews),
                "animation_mapping": any(page.scene.motion_cues for page in pages),
            },
            scanner_version=self.scanner_version,
        )
        return manifest.model_copy(update={"manifest_hash": _manifest_hash(manifest)})

    @staticmethod
    def _validate_package(path: Path, policy: FidelityPolicy) -> None:
        if not zipfile.is_zipfile(path):
            raise FidelityScanError("fidelity_package_invalid", "PPTX 不是有效的 OOXML 压缩包")
        try:
            with zipfile.ZipFile(path) as archive:
                total = 0
                for info in archive.infolist():
                    name = info.filename.replace("\\", "/")
                    if name.startswith("/") or any(part == ".." for part in name.split("/")):
                        raise FidelityScanError("fidelity_path_escape", "PPTX 包含越界路径")
                    total += info.file_size
                    if total > policy.max_xml_bytes:
                        raise FidelityScanError("fidelity_xml_limit", "PPTX 解压内容超过安全上限")
                    lowered = name.lower()
                    if any(
                        token in lowered
                        for token in ("vbaproject.bin", "activex/", "embeddings/", "externallinks/")
                    ):
                        raise FidelityScanError(
                            "fidelity_active_content", "PPTX 包含被禁止的宏、OLE 或外部链接"
                        )
        except zipfile.BadZipFile as error:
            raise FidelityScanError("fidelity_package_invalid", "PPTX 压缩包损坏") from error


def _extract_scene(
    slide: Any, index: int, page_id: UUID, slide_width: int, slide_height: int
) -> tuple[SlideScene, list[FidelityIssue]]:
    issues: list[FidelityIssue] = []
    shapes: list[SlideShape] = []
    for z_order, shape in enumerate(slide.shapes):
        try:
            x = max(0.0, min(1.0, shape.left / slide_width))
            y = max(0.0, min(1.0, shape.top / slide_height))
            width = max(0.0, min(1.0, shape.width / slide_width))
            height = max(0.0, min(1.0, shape.height / slide_height))
        except (AttributeError, TypeError, ZeroDivisionError):
            x = y = width = height = 0.0
        text = shape.text if getattr(shape, "has_text_frame", False) else ""
        kind = str(getattr(shape, "shape_type", "unknown")).lower()
        shapes.append(
            SlideShape(
                shape_id=str(getattr(shape, "shape_id", z_order)),
                name=shape.name or f"shape-{z_order}",
                kind=kind,
                z_order=z_order,
                x=x,
                y=y,
                width=width,
                height=height,
                rotation=float(getattr(shape, "rotation", 0) or 0),
                text=text[:10_000],
                text_style={"has_text": bool(text)},
            )
        )
    if not shapes:
        issues.append(
            FidelityIssue(
                code="slide_empty", message="页面未发现可渲染元素", action="检查 PPTX 页面内容"
            )
        )
    return SlideScene(
        slide_id=page_id, page_index=index, width=slide_width, height=slide_height, shapes=shapes
    ), issues


def _extract_motion(path: Path, index: int, page_id: UUID) -> list[MotionCue]:
    slide_path = f"ppt/slides/slide{index}.xml"
    try:
        with zipfile.ZipFile(path) as archive:
            xml = archive.read(slide_path).decode("utf-8", errors="replace")
    except (KeyError, OSError, UnicodeDecodeError):
        return []
    if "<p:timing" not in xml:
        return []
    cues: list[MotionCue] = []
    effects = re.findall(r"<p:animEffect[^>]*filter=\"([^\"]+)\"[^>]*>", xml)
    target_ids = re.findall(r"<p:spTgt[^>]*spid=\"([^\"]+)\"", xml)
    for sequence, effect in enumerate(effects):
        effect_lower = effect.lower()
        supported = any(name in effect_lower for name in ("fade", "wipe", "fly", "zoom", "appear"))
        cues.append(
            MotionCue(
                cue_id=uuid5(page_id, f"motion:{sequence}:{effect}"),
                shape_ids=[target_ids[sequence] if sequence < len(target_ids) else "slide"],
                sequence=sequence,
                start_ms=sequence * 500,
                duration_ms=500,
                entrance=effect_lower,
                support=MotionSupport.SUPPORTED if supported else MotionSupport.DEGRADED,
                source_effect=effect,
            )
        )
    return cues


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _manifest_hash(manifest: SlideFidelityManifest) -> str:
    payload = manifest.model_dump(mode="json", exclude={"manifest_hash"})
    return hashlib.sha256(
        json.dumps(payload, sort_keys=True, separators=(",", ":")).encode()
    ).hexdigest()
